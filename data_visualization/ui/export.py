"""
Export section: PNG, SVG, HTML, PDF, Python, Notebook, PPTX download buttons.

Kaleido (fig.to_image) launches a headless Chromium process — ~1-2s per call.
All kaleido calls are lazy: they only run when the user clicks "Generate", and
the result is cached in session_state keyed to the chart config hash so that
subsequent rerenders of the same chart are instant.

fig.to_html() and the code exports are fast and always available immediately.
"""

import io
import json
import streamlit as st
import plotly.graph_objects as go
from typing import Optional

_EXPORT_CACHE_KEY = "viz_export_cache"


def _get_export_cache() -> dict:
    return st.session_state.setdefault(_EXPORT_CACHE_KEY, {})


def _cache_bytes(cfg_hash: str, fmt: str, data: bytes) -> None:
    cache = _get_export_cache()
    entry = cache.setdefault(cfg_hash, {})
    entry[fmt] = data
    # Keep at most 5 unique chart configs in cache
    if len(cache) > 5:
        oldest = next(iter(cache))
        del cache[oldest]
    st.session_state[_EXPORT_CACHE_KEY] = cache


def _get_cached(cfg_hash: str, fmt: str) -> Optional[bytes]:
    return _get_export_cache().get(cfg_hash, {}).get(fmt)


def _lazy_image_button(
    fig: go.Figure,
    fmt: str,
    label_generate: str,
    label_download: str,
    filename: str,
    mime: str,
    cfg_hash: str,
    button_key: str,
    width: int = 1200,
    height: int = 800,
) -> None:
    """
    Show a download button if bytes are already cached, otherwise show a
    'Generate' button. One kaleido call per unique (cfg_hash, fmt) pair.
    """
    cached = _get_cached(cfg_hash, fmt)
    if cached is not None:
        st.download_button(label_download, cached, filename, mime,
                           key=f"dl_{button_key}", use_container_width=True)
    else:
        if st.button(label_generate, key=f"gen_{button_key}",
                     use_container_width=True):
            with st.spinner(f"Generating {fmt.upper()}…"):
                try:
                    data = fig.to_image(format=fmt, width=width, height=height)
                    _cache_bytes(cfg_hash, fmt, data)
                    st.rerun()
                except Exception as e:
                    st.error(f"{fmt.upper()} export failed: {e}")
                    if fmt in ("png", "svg", "pdf"):
                        st.caption("💡 Install kaleido: `pip install kaleido`")


def render_export_section(
    fig: go.Figure,
    chart_mode: str,
    chart_type: str,
    selected_table: str,
    x_col: str,
    y_col: str,
    color_col: str,
    cfg_hash: str = "",
) -> None:
    """
    Render Export Chart and Export More Formats sections.

    PNG / SVG / PDF / PPTX use lazy kaleido generation (button-triggered,
    cached per cfg_hash). HTML, Python, and Notebook are always instant.
    """
    if fig is None:
        return

    export_chart_name = chart_mode if chart_mode != "basic" else chart_type
    export_height = 800
    h = cfg_hash or str(id(fig))  # fallback key when no hash provided

    st.divider()
    st.subheader("📥 Export Chart")
    st.caption(
        "HTML is instant. PNG / SVG / PDF / PPTX are generated on first click "
        "and cached — subsequent downloads are instant."
    )

    col1, col2, col3 = st.columns(3)

    with col1:
        _lazy_image_button(
            fig, "png",
            label_generate="📥 Generate PNG",
            label_download="📥 Download PNG",
            filename=f"chart_{export_chart_name}_{selected_table}.png",
            mime="image/png",
            cfg_hash=h,
            button_key=f"png_{h[:8]}",
            height=export_height,
        )

    with col2:
        _lazy_image_button(
            fig, "svg",
            label_generate="📐 Generate SVG",
            label_download="📐 Download SVG",
            filename=f"chart_{export_chart_name}_{selected_table}.svg",
            mime="image/svg+xml",
            cfg_hash=h,
            button_key=f"svg_{h[:8]}",
            height=export_height,
        )

    with col3:
        # HTML via fig.to_html() — no kaleido, always instant
        try:
            html_str = fig.to_html(full_html=False)
            st.download_button(
                "🌐 Download HTML",
                html_str.encode(),
                f"chart_{export_chart_name}_{selected_table}.html",
                "text/html",
                key=f"dl_html_{h[:8]}",
                use_container_width=True,
            )
        except Exception as e:
            st.error(f"HTML export failed: {e}")

    st.markdown("---")
    st.subheader("📦 Export More Formats")
    st.caption("Generate reports and reusable code for your chart.")

    e1, e2, e3 = st.columns(3)

    with e1:
        _lazy_image_button(
            fig, "pdf",
            label_generate="📄 Generate PDF",
            label_download="📄 Download PDF",
            filename=f"chart_{export_chart_name}_{selected_table}.pdf",
            mime="application/pdf",
            cfg_hash=h,
            button_key=f"pdf_{h[:8]}",
            height=export_height,
        )

    with e2:
        python_script = f"""import pandas as pd
import plotly.express as px

# Load your data
# df = pd.read_csv("{selected_table}_current.csv")

fig = px.{chart_type}(
    df,
    x={repr(x_col if x_col != 'None' else None)},
    y={repr(y_col if y_col != 'None' else None)},
    color={repr(color_col if color_col != 'None' else None)},
)
fig.show()
"""
        st.download_button(
            "🐍 Download Python",
            python_script.encode(),
            f"chart_{export_chart_name}_{selected_table}.py",
            "text/x-python",
            key=f"dl_py_{h[:8]}",
            use_container_width=True,
        )

    with e3:
        notebook = {
            "cells": [
                {
                    "cell_type": "markdown",
                    "metadata": {},
                    "source": [f"# Chart Export: {export_chart_name}\n"],
                },
                {
                    "cell_type": "code",
                    "execution_count": None,
                    "metadata": {},
                    "outputs": [],
                    "source": [
                        "import pandas as pd\n",
                        "import plotly.express as px\n\n",
                        f'# df = pd.read_csv("{selected_table}_current.csv")\n',
                        "df.head()\n",
                    ],
                },
                {
                    "cell_type": "code",
                    "execution_count": None,
                    "metadata": {},
                    "outputs": [],
                    "source": [
                        f"fig = px.{chart_type}(df, "
                        f"x={repr(x_col if x_col != 'None' else None)}, "
                        f"y={repr(y_col if y_col != 'None' else None)}, "
                        f"color={repr(color_col if color_col != 'None' else None)})\n",
                        "fig.show()\n",
                    ],
                },
            ],
            "metadata": {
                "kernelspec": {
                    "display_name": "Python 3",
                    "language": "python",
                    "name": "python3",
                },
                "language_info": {"name": "python", "version": "3.x"},
            },
            "nbformat": 4,
            "nbformat_minor": 5,
        }
        st.download_button(
            "📓 Download Notebook",
            json.dumps(notebook, indent=2).encode(),
            f"chart_{export_chart_name}_{selected_table}.ipynb",
            "application/json",
            key=f"dl_nb_{h[:8]}",
            use_container_width=True,
        )

    # PPTX — reuses cached PNG bytes to avoid a second kaleido call
    ppt_col, _, _ = st.columns(3)
    with ppt_col:
        cached_png = _get_cached(h, "png")
        if cached_png is not None:
            # PNG already generated — build PPTX from cached bytes, no kaleido
            try:
                from pptx import Presentation
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                image_stream = io.BytesIO(cached_png)
                slide.shapes.add_picture(image_stream, left=0, top=0,
                                         width=prs.slide_width)
                pptx_stream = io.BytesIO()
                prs.save(pptx_stream)
                st.download_button(
                    "📊 Download PowerPoint",
                    pptx_stream.getvalue(),
                    f"chart_{export_chart_name}_{selected_table}.pptx",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"dl_pptx_{h[:8]}",
                    use_container_width=True,
                )
            except ImportError:
                st.caption("💡 Install python-pptx: `pip install python-pptx`")
            except Exception as e:
                st.error(f"PowerPoint export failed: {e}")
        else:
            st.caption("📊 PowerPoint — generate PNG first (button above).")
